from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR

prs = Presentation()

def generateslide(summary):
    list = summary.split('\n')
    for str in list:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  
        
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        listsentence = str.split('.')
        
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        textframe = textbox.text_frame
        
        listtext=['']
        
        for sentence in listsentence:
            senend=sentence+'\n'
            listtext.append(listtext[-1]+senend)
        textframe.text = listtext[-1]

s = '  % python slide.py\ (airst % rm basic_presentation.pptx\ '

generateslide(s)

prs.save('PDFtemplate.pptx')
